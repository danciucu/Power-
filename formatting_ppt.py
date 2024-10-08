from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


class Textbox():
    def __init__(self) -> None:
        pass

    def position(left, top, width, height):
        position_left = Inches(left)
        position_top = Inches(top)
        position_width = Inches(width)
        position_height = Inches(height)

        return position_left, position_top, position_width, position_height

    def create_textbox(position, slide, text, size, is_bold):
        # add textbox
        textbox = slide.shapes.add_textbox(*position)
        text_frame = textbox.text_frame
        # set the text in the textbox
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(size)
        paragraph.font.bold = is_bold
        paragraph.alignment = PP_ALIGN.CENTER