import os
from pptx import Presentation

def extract_images_and_captions_from_ppt(pptx_path, output_folder):
    # Load the presentation
    prs = Presentation(pptx_path)
    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.mkdir(output_folder)
    # Initialize image counter
    image_counter = 1
    # Array to hold image file paths
    image_path_array = []
    # Array to hold captions
    caption_array = []
    # Array to hold image file paths and captions
    image_caption_dictionary = {}

    # Loop through slides
    for slide_num, slide in enumerate(prs.slides):
        # Loop through shapes in each slide
        for shape in slide.shapes:
            # Check if the shape contains a picture
            if hasattr(shape, "image"):
                # Get the image in the shape
                image = shape.image
                # Define output image name
                image_name = f"image_{slide_num + 1}.jpg"
                # Define output image path
                image_path = os.path.join(output_folder, image_name)
                
                # Save the image
                with open(image_path, "wb") as img_file:
                    img_file.write(image.blob)

                image_path_array.append(image_path)

            if hasattr(shape, 'text') and ("Inspection" not in shape.text) and (len(shape.text) > 1):
                caption_array.append(str(shape.text))

    #print(caption_array)

    for i in range(len(image_path_array)):
        image_caption_dictionary[image_path_array[i]] = caption_array[i]

    return image_caption_dictionary


def extract_inspectors_from_ppt(pptx_path):
    # Load the presentation
    prs = Presentation(pptx_path)

    # Default value if no matching text is found
    inspectors = 'Check Who Inspected The Bridge'
    
    # Loop through slides
    for slide_num, slide in enumerate(prs.slides):
        # Check only the first slide
        if slide_num == 0:
            # Loop through shapes in the first slide
            for shape in slide.shapes:
                if hasattr(shape, 'text') and ("Inspection" in shape.text):
                    inspectors = str(shape.text)
                    break
                
    return inspectors